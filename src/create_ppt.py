# @title create_ppt.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def create_presentation():
    prs = Presentation()

    # 색상 테마 (진한 파스텔톤 - 사용자 선호)
    # Deep Pastel Green (맑은 날씨/긍정), Deep Pastel Gray (기본 텍스트), Deep Pastel Red (강조)
    COLOR_PRIMARY = RGBColor(77, 166, 105)  # 진한 파스텔 초록
    COLOR_SECONDARY = RGBColor(90, 90, 90)  # 진한 회색
    COLOR_ACCENT = RGBColor(235, 110, 90)  # 진한 파스텔 레드 (사과색)

    def add_slide(title_text, content_text_list):
        slide_layout = prs.slide_layouts[1]  # 제목 + 내용 레이아웃
        slide = prs.slides.add_slide(slide_layout)

        # 제목 설정
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
        title.text_frame.paragraphs[0].font.bold = True

        # 내용 설정
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()  # 기존 텍스트 제거

        for i, text in enumerate(content_text_list):
            p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(20)
            p.font.color.rgb = COLOR_SECONDARY
            p.space_after = Pt(14)

            # 첫 번째 줄은 그대로, 그 다음부터는 들여쓰기 등 조정 가능
            if text.startswith("-"):
                p.level = 1

    # 1. 표지
    slide_layout = prs.slide_layouts[0]  # 제목 슬라이드
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.shapes.placeholders[1]

    title.text = "입찰메이트 RAG 시스템 개발\n중간 보고서"
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    title.text_frame.paragraphs[0].font.bold = True

    subtitle.text = "발표자: 개발팀\n2025년 12월 22일"
    subtitle.text_frame.paragraphs[0].font.color.rgb = COLOR_SECONDARY

    # 2. 개발 배경
    add_slide("1. 개발 배경 및 목표", [
        "배경:",
        "- 공공 입찰 제안요청서(RFP)의 방대한 분량",
        "- 수동 검토로 인한 시간 소요 및 휴먼 에러 발생",
        "목표:",
        "- LLM 기반 RAG 시스템 구축",
        "- 사용자의 자연어 질문에 대한 정확한 정보 추출 및 답변"
    ])

    # 3. 시스템 아키텍처
    add_slide("2. 시스템 아키텍처", [
        "Data Pipeline:",
        "- CSV 및 HWP/PDF 비정형 데이터 로딩 및 전처리",
        "Vector DB:",
        "- ChromaDB 활용, 문서 임베딩 저장 및 검색",
        "Generator:",
        "- LangChain & GPT-5 활용",
        "User Interface:",
        "- Streamlit 기반의 대화형 웹 애플리케이션"
    ])

    # 4. 핵심 기능
    add_slide("3. 핵심 구현 기능", [
        "💬 대화형 챗봇 Interface",
        "- 입찰 관련 질문에 실시간 답변 제공",
        "🛠️ 관리자 평가 도구 (Admin Sidebar)",
        "- RAG 모델의 정확도 실시간 측정 기능 탑재",
        "- 정답셋(Ground Truth) 비교 알고리즘 적용"
    ])

    # 5. 성능 평가 및 이슈
    add_slide("4. 초기 성능 평가 및 이슈", [
        "평가 현황:",
        "- 총 300개 테스트 문항 기반 평가 진행",
        "- 초기 정확도: 목표치 미달 (약 60% 대)",
        "발견된 문제점:",
        "- Hallucination: 유사한 타 사업/연도 문서를 참조하는 오류",
        "- Format Mismatch: 파일 확장자(hwp) 질문에 문서 성격(RFP)으로 오답 처리"
    ])

    # 6. 문제 해결 (Troubleshooting)
    add_slide("5. 트러블 슈팅 (개선 방안)", [
        "✅ 데이터 로더 개선 (Data Loader)",
        "- 파일명에서 확장자(.hwp, .pdf) 추출 후 메타데이터 주입",
        "✅ 프롬프트 엔지니어링 강화",
        "- '질문의 사업명과 일치하는 정보만 참조'하도록 제약 조건 추가",
        "- 문서 포맷 질문 시 확장자를 포함하도록 지시",
        "👉 결과: 오답률 대폭 감소 및 신뢰도 향상 기대"
    ])

    # 7. 향후 계획
    add_slide("6. 향후 계획 (Roadmap)", [
        "📊 데이터 시각화 (Dashboard)",
        "- 발주 기관별, 예산별 통계 그래프 구현",
        "- 직관적인 파스텔톤 컬러 차트 적용 예정",
        "🚀 시스템 고도화",
        "- 최종 정확도 95% 이상 달성 목표",
        "- UI/UX 개선 및 배포 준비"
    ])

    # 저장
    save_name = "midterm_presentation_rag.pptx"
    prs.save(save_name)
    print(f"✅ PPT 파일 생성 완료: {save_name}")


if __name__ == "__main__":
    create_presentation()